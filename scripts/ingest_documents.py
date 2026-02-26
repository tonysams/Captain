"""
Document ingestion script for the qualitative research pipeline.

Reads PDF, DOCX, PPTX, MD, and Pages files from a source directory and
extracts their text content to output/clean/ as .txt or .md files.

Usage:
    python3 scripts/ingest_documents.py \
        --source-dir "/Users/tony/Library/Mobile Documents/com~apple~CloudDocs/CoWork/AI_Symposium" \
        --output-dir output/clean/
"""

import argparse
import json
import shutil
import subprocess
import sys
from pathlib import Path


def ingest_markdown(path: Path, output_dir: Path) -> dict:
    """Copy markdown files as-is — they're already clean text."""
    out_path = output_dir / path.name
    shutil.copy2(path, out_path)
    word_count = len(path.read_text(encoding="utf-8", errors="replace").split())
    return {"file": str(path), "type": "markdown", "output": str(out_path), "words": word_count}


def ingest_pdf(path: Path, output_dir: Path) -> dict:
    """Extract text from PDF using pdfplumber (more reliable than pypdf for mixed content)."""
    try:
        import pdfplumber
        text_parts = []
        with pdfplumber.open(path) as pdf:
            for page in pdf.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
        full_text = "\n\n".join(text_parts)
    except Exception:
        # Fallback to pypdf
        try:
            import pypdf
            reader = pypdf.PdfReader(str(path))
            text_parts = []
            for page in reader.pages:
                text = page.extract_text()
                if text:
                    text_parts.append(text)
            full_text = "\n\n".join(text_parts)
        except Exception as e:
            return {"file": str(path), "type": "pdf", "error": str(e), "skipped": True}

    if not full_text.strip():
        return {"file": str(path), "type": "pdf", "error": "No text extracted (may be scanned image)", "skipped": True}

    stem = path.stem.replace(" ", "_")
    out_path = output_dir / f"{stem}.txt"
    out_path.write_text(full_text, encoding="utf-8")
    return {"file": str(path), "type": "pdf", "output": str(out_path), "words": len(full_text.split())}


def ingest_docx(path: Path, output_dir: Path) -> dict:
    """Extract text from .docx using python-docx."""
    try:
        import docx
        doc = docx.Document(str(path))
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        # Also extract table text
        for table in doc.tables:
            for row in table.rows:
                row_text = "\t".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    paragraphs.append(row_text)
        full_text = "\n\n".join(paragraphs)
    except Exception as e:
        return {"file": str(path), "type": "docx", "error": str(e), "skipped": True}

    stem = path.stem.replace(" ", "_")
    out_path = output_dir / f"{stem}.txt"
    out_path.write_text(full_text, encoding="utf-8")
    return {"file": str(path), "type": "docx", "output": str(out_path), "words": len(full_text.split())}


def ingest_pptx(path: Path, output_dir: Path) -> dict:
    """Extract text from .pptx using python-pptx."""
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slide_texts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            if texts:
                slide_texts.append(f"[Slide {i}]\n" + "\n".join(texts))
        full_text = "\n\n".join(slide_texts)
    except Exception as e:
        return {"file": str(path), "type": "pptx", "error": str(e), "skipped": True}

    stem = path.stem.replace(" ", "_")
    out_path = output_dir / f"{stem}.txt"
    out_path.write_text(full_text, encoding="utf-8")
    return {"file": str(path), "type": "pptx", "output": str(out_path), "words": len(full_text.split())}


def ingest_pages(path: Path, output_dir: Path) -> dict:
    """Extract text from .pages using macOS textutil (built-in)."""
    stem = path.stem.replace(" ", "_")
    out_path = output_dir / f"{stem}.txt"
    try:
        result = subprocess.run(
            ["textutil", "-convert", "txt", "-stdout", str(path)],
            capture_output=True, text=True, timeout=30
        )
        if result.returncode != 0 or not result.stdout.strip():
            # Try alternate approach: convert to txt file
            result2 = subprocess.run(
                ["textutil", "-convert", "txt", "-output", str(out_path), str(path)],
                capture_output=True, text=True, timeout=30
            )
            if result2.returncode != 0:
                return {"file": str(path), "type": "pages", "error": result2.stderr, "skipped": True}
            full_text = out_path.read_text(encoding="utf-8", errors="replace")
        else:
            full_text = result.stdout
            out_path.write_text(full_text, encoding="utf-8")
    except Exception as e:
        return {"file": str(path), "type": "pages", "error": str(e), "skipped": True}

    return {"file": str(path), "type": "pages", "output": str(out_path), "words": len(full_text.split())}


HANDLERS = {
    ".md": ingest_markdown,
    ".pdf": ingest_pdf,
    ".docx": ingest_docx,
    ".pptx": ingest_pptx,
    ".pages": ingest_pages,
}


def walk_source(source_dir: Path) -> list:
    """Recursively find all supported document files."""
    files = []
    for path in sorted(source_dir.rglob("*")):
        if path.is_file() and path.suffix.lower() in HANDLERS:
            if not any(part.startswith(".") for part in path.parts):
                files.append(path)
    return files


def main():
    parser = argparse.ArgumentParser()
    parser.add_argument("--source-dir", required=True, help="Directory containing research documents")
    parser.add_argument("--output-dir", required=True, help="Directory to write extracted text files")
    args = parser.parse_args()

    source_dir = Path(args.source_dir)
    output_dir = Path(args.output_dir)
    output_dir.mkdir(parents=True, exist_ok=True)

    if not source_dir.exists():
        print(f"Error: source directory does not exist: {source_dir}", file=sys.stderr)
        sys.exit(1)

    files = walk_source(source_dir)
    if not files:
        print(f"No supported documents found in {source_dir}", file=sys.stderr)
        sys.exit(1)

    print(f"Found {len(files)} documents to ingest from {source_dir}")
    results = []

    for path in files:
        handler = HANDLERS[path.suffix.lower()]
        print(f"  Ingesting [{path.suffix}]: {path.name}")
        result = handler(path, output_dir)
        results.append(result)
        if result.get("skipped"):
            print(f"    SKIPPED: {result.get('error', 'unknown error')}")
        else:
            print(f"    OK: {result.get('words', 0):,} words → {Path(result['output']).name}")

    ingested = [r for r in results if not r.get("skipped")]
    skipped = [r for r in results if r.get("skipped")]

    print(f"\nIngestion complete: {len(ingested)} files processed, {len(skipped)} skipped")
    print(json.dumps(results, indent=2))


if __name__ == "__main__":
    main()
